#!/usr/bin/env python3
"""
extract_quadruplets.py
==============================
Извлекает квадруплеты (s, r, o, t) из документов (PDF, DOCX, PPTX, TXT)
с помощью DeepSeek, нормализует короткие имена через Natasha,
сохраняет результат в JSONL (extract_data/quadruplets.jsonl).

Зависимости:
    pip install openai pymupdf python-docx python-pptx tqdm natasha

Использование:
    python extract_quadruplets.py --input_dir new_docs
    python extract_quadruplets.py --input_dir new_docs --workers 30
    python extract_quadruplets.py --input_dir new_docs --reprocess
    python extract_quadruplets.py --input_dir new_docs --output my_output.jsonl
    python extract_quadruplets.py --input_dir new_docs --model deepseek-chat --api_key sk-...
    python extract_quadruplets.py --input_dir new_docs --mark-done   # пометить без LLM-вызовов
    python extract_quadruplets.py --input_dir new_docs --debug       # печатать каждый квадруплет
"""

import os
import re
import sys
import json
import time
import hashlib
import argparse
import logging
import threading
from pathlib import Path
from typing import Optional
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime, timezone

from tqdm import tqdm

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None  # type: ignore[assignment]

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None  # type: ignore[assignment]

try:
    from pptx import Presentation
except ImportError:
    Presentation = None  # type: ignore[assignment]

try:
    from openai import OpenAI
except ImportError:
    OpenAI = None  # type: ignore[assignment]

try:
    from natasha import Segmenter, MorphVocab, NewsEmbedding, NewsMorphTagger, Doc as NatashaDoc
    _natasha_available = True
except ImportError:
    _natasha_available = False

logging.basicConfig(level=logging.INFO, format="%(levelname)s | %(message)s")
log = logging.getLogger(__name__)

# ─────────────────────────────────────────────
# СТАТИСТИКА ТОКЕНОВ
# ─────────────────────────────────────────────

_token_stats: dict[str, int] = {"prompt": 0, "completion": 0, "total": 0, "calls": 0}
_token_lock = threading.Lock()


def _add_tokens(prompt: int = 0, completion: int = 0) -> None:
    with _token_lock:
        _token_stats["prompt"]     += prompt
        _token_stats["completion"] += completion
        _token_stats["total"]      += prompt + completion
        _token_stats["calls"]      += 1


# ─────────────────────────────────────────────
# 1. ЧТЕНИЕ ДОКУМЕНТОВ
# ─────────────────────────────────────────────

def read_pdf(path: str) -> str:
    with fitz.open(path) as doc:
        return "\n".join(page.get_text() for page in doc)

def read_docx(path: str) -> str:
    doc = DocxDocument(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

def read_pptx(path: str) -> str:
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n".join(parts)

READERS = {
    ".pdf":  read_pdf,
    ".docx": read_docx,
    ".pptx": read_pptx,
    ".txt":  lambda p: Path(p).read_text(encoding="utf-8", errors="ignore"),
}

def read_document(path: str) -> Optional[str]:
    ext = Path(path).suffix.lower()
    reader = READERS.get(ext)
    if reader is None:
        log.warning(f"Unsupported format: {path}")
        return None
    try:
        text = reader(path)
        log.info(f"Read {path} → {len(text):,} chars")
        return text
    except Exception as e:
        log.error(f"Cannot read {path}: {e}")
        return None

def list_documents(input_dir: str) -> list[str]:
    paths = [str(p) for p in Path(input_dir).rglob("*")
             if p.suffix.lower() in READERS and p.is_file()]
    log.info(f"Found {len(paths)} documents in '{input_dir}'")
    return paths


# ─────────────────────────────────────────────
# 2. ПУТИ И ФАЙЛЫ ДАННЫХ
# ─────────────────────────────────────────────

OUTPUT_DIR              = Path("extract_data")
CUSTOM_ID_REGISTRY_FILE = str(OUTPUT_DIR / "custom_ids_registry.json")
PROCESSED_MANIFEST_FILE = str(OUTPUT_DIR / "processed_files.json")


# ─────────────────────────────────────────────
# 3. РЕЕСТР КАСТОМНЫХ ID
# ─────────────────────────────────────────────

_custom_registry: dict = {
    "entities":        {},
    "properties":      {},
    "next_entity_seq": 1,
    "next_prop_seq":   1,
}
_registry_lock = threading.Lock()


def _load_custom_registry() -> None:
    global _custom_registry
    if os.path.exists(CUSTOM_ID_REGISTRY_FILE):
        with open(CUSTOM_ID_REGISTRY_FILE, encoding="utf-8") as f:
            _custom_registry = json.load(f)
        log.info(
            f"Custom ID registry: {len(_custom_registry['entities'])} entities, "
            f"{len(_custom_registry['properties'])} properties"
        )

def _save_custom_registry() -> None:
    tmp = CUSTOM_ID_REGISTRY_FILE + ".tmp"
    with open(tmp, "w", encoding="utf-8") as f:
        json.dump(_custom_registry, f, ensure_ascii=False)
    os.replace(tmp, CUSTOM_ID_REGISTRY_FILE)

def get_or_create_entity_id(name: str) -> str:
    key = name.strip().lower()
    with _registry_lock:
        if key in _custom_registry["entities"]:
            return _custom_registry["entities"][key]
        seq    = _custom_registry["next_entity_seq"]
        new_id = f"CQ_{seq:06d}"
        _custom_registry["entities"][key] = new_id
        _custom_registry["next_entity_seq"] = seq + 1
    return new_id

def get_or_create_property_id(name: str) -> str:
    key = name.strip().lower()
    with _registry_lock:
        if key in _custom_registry["properties"]:
            return _custom_registry["properties"][key]
        seq    = _custom_registry["next_prop_seq"]
        new_id = f"CP_{seq:06d}"
        _custom_registry["properties"][key] = new_id
        _custom_registry["next_prop_seq"] = seq + 1
    return new_id


# ─────────────────────────────────────────────
# 4. МАНИФЕСТ ОБРАБОТАННЫХ ФАЙЛОВ
# ─────────────────────────────────────────────

_processed_manifest: dict[str, dict] = {}


def _file_sha256(filepath: str) -> str:
    h = hashlib.sha256()
    with open(filepath, "rb") as f:
        while chunk := f.read(65536):
            h.update(chunk)
    return h.hexdigest()

def _load_processed_manifest() -> None:
    global _processed_manifest
    if os.path.exists(PROCESSED_MANIFEST_FILE):
        with open(PROCESSED_MANIFEST_FILE, encoding="utf-8") as f:
            _processed_manifest = json.load(f)
        log.info(f"Manifest: {len(_processed_manifest)} already-processed file(s)")

def _save_processed_manifest() -> None:
    tmp = PROCESSED_MANIFEST_FILE + ".tmp"
    with open(tmp, "w", encoding="utf-8") as f:
        json.dump(_processed_manifest, f, ensure_ascii=False, indent=2)
    os.replace(tmp, PROCESSED_MANIFEST_FILE)

def _mark_processed(filepath: str, quads_count: int, sha: str) -> None:
    key = str(Path(filepath).resolve())
    _processed_manifest[key] = {
        "hash":         sha,
        "processed_at": datetime.now(timezone.utc).isoformat(timespec="seconds"),
        "quads":        quads_count,
    }
    _save_processed_manifest()


# ─────────────────────────────────────────────
# 5. DeepSeek LLM
# ─────────────────────────────────────────────

DEEPSEEK_BASE_URL      = "https://api.deepseek.com"
DEEPSEEK_DEFAULT_MODEL = "deepseek-chat"
DEEPSEEK_API_KEY = os.environ.get("DEEPSEEK_API_KEY", "")

EXTRACTION_SYSTEM = """You are a knowledge graph construction expert.
Extract factual quadruplets (subject, relation, object, time) from the text.

Rules:
1. Subject and Object must be named entities (persons, organizations, places, works, concepts).
2. Relation must be a short descriptive label (e.g. "employer", "member of", "has part", "regulates", "approved by").
3. Time: use years only. If exact year unknown, omit. Format: {"start": "YYYY", "end": "YYYY"}.
   If a single year, set start == end. If open-ended (still active), use current year as end.
4. Output ONLY a valid JSON object with a single key "quadruplets" containing the array.
5. Language: Relations MUST be in Russian since source texts are in Russian. Use short Russian phrases (e.g. "является директором", "входит в состав", "чистая прибыль составила"). Subject and Object names: preserve original form from the text.

Output format (strict JSON):
{
  "quadruplets": [
    {"s": "Subject Name", "r": "relation label", "o": "Object Name", "t_start": "YYYY", "t_end": "YYYY"}
  ]
}

If no quadruplets found, return: {"quadruplets": []}
"""

def chunk_text(text: str, max_chars: int = 6000) -> list[str]:
    chunks, step = [], max(1, max_chars - 500)
    for i in range(0, len(text), step):
        chunks.append(text[i:i + max_chars])
    return chunks

def parse_llm_json(raw: str) -> list[dict]:
    raw = re.sub(r"```(?:json)?", "", raw.strip()).replace("```", "").strip()
    try:
        parsed = json.loads(raw)
        if isinstance(parsed, list):
            return parsed
        for key in ("quadruplets", "data", "results"):
            if key in parsed and isinstance(parsed[key], list):
                return parsed[key]
    except json.JSONDecodeError:
        pass
    idx = raw.find("[")
    if idx == -1:
        log.warning(f"No JSON array found in LLM response: {raw[:200]}")
        return []
    try:
        result, _ = json.JSONDecoder().raw_decode(raw, idx)
        if isinstance(result, list):
            return result
    except json.JSONDecodeError as e:
        log.warning(f"JSON parse error: {e}\nRaw: {raw[:300]}")
    return []


def _process_chunk(chunk: str, client, model: str,
                   max_retries: int = 5) -> tuple[list[dict], int, int]:
    for attempt in range(max_retries):
        try:
            resp = client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": EXTRACTION_SYSTEM},
                    {"role": "user",   "content": f"Extract quadruplets from this text:\n\n{chunk}"},
                ],
                temperature=0,
                response_format={"type": "json_object"},
            )
            raw   = resp.choices[0].message.content
            usage = resp.usage
            return (parse_llm_json(raw),
                    usage.prompt_tokens if usage else 0,
                    usage.completion_tokens if usage else 0)
        except Exception as e:
            status = getattr(getattr(e, "response", None), "status_code", None)
            if status in (429, 500, 502, 503) and attempt < max_retries - 1:
                wait = min(2 ** attempt, 30)
                log.warning(f"API {status} — retry {attempt+1}/{max_retries} after {wait}s")
                time.sleep(wait)
            else:
                raise
    raise RuntimeError("unreachable")


def extract_with_deepseek(text: str, api_key: str, model: str,
                          workers: int = 50) -> tuple[list[dict], int]:
    if OpenAI is None:
        log.error("Install openai: pip install openai")
        return [], 0

    client = OpenAI(api_key=api_key, base_url=DEEPSEEK_BASE_URL, timeout=60.0)
    chunks = chunk_text(text)
    all_quads: list[dict] = []
    failed = 0

    with ThreadPoolExecutor(max_workers=min(workers, len(chunks))) as executor:
        futures = {executor.submit(_process_chunk, chunk, client, model): i
                   for i, chunk in enumerate(chunks)}
        bar = tqdm(as_completed(futures), total=len(chunks),
                   desc="  chunks", unit="chunk", leave=False)
        for future in bar:
            try:
                quads, pt, ct = future.result()
                all_quads.extend(quads)
                _add_tokens(pt, ct)
                bar.set_postfix(tokens=f"{_token_stats['total']:,}")
            except Exception as e:
                log.error(f"Chunk error: {e}")
                failed += 1

    if failed:
        log.warning(f"extract_with_deepseek: {failed}/{len(chunks)} chunks failed")
    return all_quads, failed


# ─────────────────────────────────────────────
# 6. НОРМАЛИЗАЦИЯ (только Natasha для коротких имён)
# ─────────────────────────────────────────────

_natasha_tools: Optional[tuple] = None
_natasha_lock  = threading.Lock()


def _get_natasha():
    global _natasha_tools
    if _natasha_tools is not None:
        return _natasha_tools
    with _natasha_lock:
        if _natasha_tools is not None:
            return _natasha_tools
        if not _natasha_available:
            log.warning("Natasha не установлена (pip install natasha) — нормализация пропущена")
            _natasha_tools = ()
        else:
            seg    = Segmenter()
            vocab  = MorphVocab()
            emb    = NewsEmbedding()
            tagger = NewsMorphTagger(emb)
            _natasha_tools = (seg, vocab, tagger)
            log.info("Natasha: инициализирована")
    return _natasha_tools


def _natasha_lemma(name: str) -> str:
    tools = _get_natasha()
    if not tools:
        return name.lower()
    try:
        seg, vocab, tagger = tools
        doc = NatashaDoc(name)
        doc.segment(seg)
        doc.tag_morph(tagger)
        for token in doc.tokens:
            token.lemmatize(vocab)
        return " ".join(t.lemma for t in doc.tokens).lower()
    except Exception:
        return name.lower()


def _natasha_lemmas_sequential(names: list[str]) -> dict[str, str]:
    return {name: _natasha_lemma(name) for name in names}


def normalize_names(names: list[str]) -> dict[str, str]:
    """
    Нормализует короткие имена (≤2 слова) через Natasha лемматизацию.
    Длинные имена (>2 слов) нормализуются только по strip+lower для точных совпадений.
    Возвращает {original_name: canonical_name}.
    """
    unique: list[str] = sorted(
        {n.strip() for n in names if len(n.strip()) >= 2},
        key=len, reverse=True,  # длинные → кандидаты в canonical
    )
    if not unique:
        return {}

    # Разбиваем на короткие (лемматизируем) и длинные (группируем по lower)
    short = [n for n in unique if len(n.split()) <= 2]
    long_ = [n for n in unique if len(n.split()) > 2]

    # Лемматизируем короткие
    lemma_map = _natasha_lemmas_sequential(short) if short else {}

    # Группируем короткие по лемме
    lemma_groups: dict[str, list[str]] = {}
    for name in short:
        key = lemma_map.get(name, name.lower())
        lemma_groups.setdefault(key, []).append(name)

    mapping: dict[str, str] = {}
    for group in lemma_groups.values():
        canon = max(group, key=len)
        for n in group:
            mapping[n] = canon

    # Длинные: группируем только по lower (точные совпадения после strip)
    long_groups: dict[str, list[str]] = {}
    for name in long_:
        long_groups.setdefault(name.lower(), []).append(name)

    for group in long_groups.values():
        canon = max(group, key=len)
        for n in group:
            mapping[n] = canon

    merged = sum(1 for v, c in mapping.items() if v != c)
    print(f"  Нормализация: {merged} вариантов смержено из {len(mapping)}")
    return mapping


# ─────────────────────────────────────────────
# 7. ПОСТРОЕНИЕ КВАДРУПЛЕТОВ
# ─────────────────────────────────────────────

def build_quadruplet(raw: dict, norm_map: dict[str, str]) -> Optional[dict]:
    s_name  = str(raw.get("s", "")).strip()
    r_name  = str(raw.get("r", "")).strip()
    o_name  = str(raw.get("o", "")).strip()
    t_start = str(raw.get("t_start", "")).strip()
    t_end   = str(raw.get("t_end", t_start)).strip()

    if not s_name or not r_name or not o_name:
        return None

    s_name = norm_map.get(s_name, s_name)
    o_name = norm_map.get(o_name, o_name)
    r_name = norm_map.get(r_name, r_name)

    s_id = get_or_create_entity_id(s_name)
    o_id = get_or_create_entity_id(o_name)
    r_id = get_or_create_property_id(r_name)

    t_block = (
        {"name": f"{t_start} - {t_end}", "type": "time",
         "prop": {"start": t_start, "end": t_end}}
        if t_start else
        {"name": "unknown", "type": "time", "prop": {}}
    )

    return {
        "s": {"name": s_name, "id": s_id},
        "r": {"name": r_name, "id": r_id, "type": "simple"},
        "o": {"name": o_name, "id": o_id},
        "t": t_block,
    }


# ─────────────────────────────────────────────
# 8. ГЛАВНАЯ ФУНКЦИЯ
# ─────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Extract knowledge graph quadruplets from documents using DeepSeek")
    parser.add_argument("--input_dir", default="new_docs")
    parser.add_argument("--output",    default=str(OUTPUT_DIR / "quadruplets.jsonl"))
    parser.add_argument("--model",     default=DEEPSEEK_DEFAULT_MODEL)
    parser.add_argument("--api_key",   default=DEEPSEEK_API_KEY)
    parser.add_argument("--workers",   type=int, default=50,
                        help="Parallel chunk workers (default: 50)")
    parser.add_argument("--reprocess", action="store_true",
                        help="Reprocess already-processed files")
    parser.add_argument("--mark-done", action="store_true",
                        help="Mark all files as processed without LLM calls")
    parser.add_argument("--debug",     action="store_true",
                        help="Debug mode: print each quadruplet")
    args = parser.parse_args()

    logging.getLogger().setLevel(logging.DEBUG if args.debug else logging.WARNING)

    # ── Режим пометки ─────────────────────────────────────
    if args.mark_done:
        OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
        _load_processed_manifest()
        paths = list_documents(args.input_dir)
        if not paths:
            print("Документы не найдены.")
            sys.exit(0)
        for filepath in tqdm(paths, desc="Marking"):
            _mark_processed(filepath, quads_count=-1, sha=_file_sha256(filepath))
        print(f"Помечено {len(paths)} файлов → {PROCESSED_MANIFEST_FILE}")
        sys.exit(0)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    _load_custom_registry()
    _load_processed_manifest()

    if not os.path.exists(args.input_dir):
        log.error(f"Input directory '{args.input_dir}' not found.")
        sys.exit(1)

    all_paths = list_documents(args.input_dir)
    if not all_paths:
        print("Документы не найдены.")
        sys.exit(0)

    # Фильтруем по манифесту.
    # SHA256 для уже известных файлов вычисляем сразу (проверка изменений),
    # для новых файлов — откладываем до момента обработки.
    to_process: list[tuple[str, str | None]] = []
    skipped = 0
    for fp in all_paths:
        key = str(Path(fp).resolve())
        if not args.reprocess and key in _processed_manifest:
            sha = _file_sha256(fp)
            if _processed_manifest[key].get("hash") == sha:
                skipped += 1
            else:
                to_process.append((fp, sha))   # файл изменился
        else:
            to_process.append((fp, None))      # новый файл, sha отложен

    print(f"Найдено документов: {len(all_paths)}  |  "
          f"К обработке: {len(to_process)}  |  Пропущено: {skipped}")

    if not to_process:
        print("Все файлы уже обработаны. Используйте --reprocess для повторной обработки.")
        sys.exit(0)

    out_path = Path(args.output)

    # При --reprocess очищаем JSONL и начинаем seen_keys с нуля
    if args.reprocess:
        out_path.open("w").close()
        seen_keys: set[tuple] = set()
    else:
        # Загружаем ключи уже сохранённых квадруплетов для кросс-сессионной дедупликации
        seen_keys = set()
        if out_path.exists():
            with open(out_path, encoding="utf-8") as f:
                for line in f:
                    line = line.strip()
                    if not line:
                        continue
                    try:
                        q = json.loads(line)
                        seen_keys.add((q["s"]["id"], q["r"]["id"], q["o"]["id"],
                                       q["t"]["prop"].get("start"), q["t"]["prop"].get("end")))
                    except Exception:
                        pass
            if seen_keys:
                log.info(f"Loaded {len(seen_keys)} existing quadruplet keys from {out_path}")

    # ── Основной цикл ─────────────────────────────────────
    total_written = 0
    start_total   = time.monotonic()

    doc_bar = tqdm(to_process, desc="Документы", unit="doc")
    for doc_idx, (filepath, doc_sha) in enumerate(doc_bar):
        doc_start = time.monotonic()
        doc_name  = Path(filepath).name
        doc_bar.set_description(f"[{doc_idx+1}/{len(to_process)}] {doc_name[:40]}")

        if doc_sha is None:
            doc_sha = _file_sha256(filepath)

        text = read_document(filepath)
        if not text:
            continue

        # Извлечение квадруплетов через LLM
        raw_quads, _ = extract_with_deepseek(
            text, api_key=args.api_key, model=args.model, workers=args.workers,
        )

        elapsed_doc = time.monotonic() - doc_start
        print(f"  {doc_name}: извлечено {len(raw_quads)} квадруплетов за {elapsed_doc:.1f}s"
              f"  (токены: {_token_stats['total']:,})")

        # Нормализация через Natasha (только короткие имена)
        all_names = (
            [str(q.get("s", "")) for q in raw_quads] +
            [str(q.get("o", "")) for q in raw_quads] +
            [str(q.get("r", "")) for q in raw_quads]
        )
        norm_map = normalize_names(all_names)

        # Построение и дедупликация квадруплетов
        doc_quads: list[dict] = []
        for rq in raw_quads:
            quad = build_quadruplet(rq, norm_map)
            if quad:
                key = (quad["s"]["id"], quad["r"]["id"], quad["o"]["id"],
                       quad["t"]["prop"].get("start"), quad["t"]["prop"].get("end"))
                if key not in seen_keys:
                    seen_keys.add(key)
                    doc_quads.append(quad)
                    if args.debug:
                        print(json.dumps(quad, ensure_ascii=False))

        # Сразу пишем в JSONL (append) — не накапливаем в памяти
        with open(out_path, "a", encoding="utf-8") as f:
            for q in doc_quads:
                f.write(json.dumps(q, ensure_ascii=False) + "\n")

        total_written += len(doc_quads)
        _save_custom_registry()
        _mark_processed(filepath, len(doc_quads), sha=doc_sha)

        # ETA
        elapsed = time.monotonic() - start_total
        done, left = doc_idx + 1, len(to_process) - doc_idx - 1
        if left > 0:
            eta = int(elapsed / done * left)
            h, rem = divmod(eta, 3600)
            m, s   = divmod(rem, 60)
            print(f"  ETA: осталось {left} doc(s) — ~{h:02d}:{m:02d}:{s:02d}")

    # ── Итог ──────────────────────────────────────────────
    total_sec = int(time.monotonic() - start_total)
    h, rem = divmod(total_sec, 3600)
    m, s   = divmod(rem, 60)
    print(f"\nГотово: добавлено {total_written} уникальных квадруплетов  |  "
          f"Токены: {_token_stats['total']:,}  |  Время: {h:02d}:{m:02d}:{s:02d}")
    log.info(f"Saved JSONL → {out_path}")

    print(f"\nDone: {total_written} new quadruplets → {out_path}")


if __name__ == "__main__":
    main()


# ─────────────────────────────────────────────
# RESET REGISTRY (call after /clear)
# ─────────────────────────────────────────────

def reset_registry() -> None:
    """Clear in-memory CQ/CP ID registry, reset counters, and persist to disk.
    Call this after /clear to ensure fresh IDs for a rebuilt graph."""
    global _custom_registry
    with _registry_lock:
        _custom_registry = {
            "entities":        {},
            "properties":      {},
            "next_entity_seq": 1,
            "next_prop_seq":   1,
        }
        _save_custom_registry()
    log.info("[reset_registry] CQ/CP ID registry cleared and saved to disk.")


# ─────────────────────────────────────────────
# BACKWARD-COMPAT SHIMS (for doc_ingestion_service.py)
# ─────────────────────────────────────────────

def collect_documents(input_dir: str) -> list:
    """Legacy API: returns [(filepath, text), ...]. New code uses list_documents()."""
    paths = list_documents(input_dir)
    result = []
    for p in paths:
        text = read_document(p)
        if text:
            result.append((p, text))
    return result


def fetch_all_wikidata_properties(use_cache: bool = True) -> dict:
    """No-op: Wikidata lookup replaced by Natasha normalization in new pipeline."""
    return {}


def _is_processed(filepath: str) -> bool:
    """Check manifest for already-processed files."""
    key = str(Path(filepath).resolve())
    return key in _processed_manifest


def deduplicate(quads: list) -> list:
    """Deduplicate quadruplets by (s_id, r_id, o_id, t_start, t_end)."""
    seen: set = set()
    result = []
    for q in quads:
        key = (
            q.get("s", {}).get("id"), q.get("r", {}).get("id"),
            q.get("o", {}).get("id"),
            q.get("t", {}).get("prop", {}).get("start"),
            q.get("t", {}).get("prop", {}).get("end"),
        )
        if key not in seen:
            seen.add(key)
            result.append(q)
    return result


def _load_prop_llm_cache() -> None:
    """No-op: Wikidata property LLM cache no longer used."""
    pass


def extract_with_ollama(text: str, model: str = "qwen3:8b") -> list:
    """Legacy shim — Ollama not implemented in new pipeline. Returns empty list."""
    log.warning("extract_with_ollama: Ollama backend removed. Use extract_with_deepseek().")
    return []


def extract_with_openai(text: str, api_key: str, model: str = "gpt-4o",
                        base_url: str = None, workers: int = 4) -> tuple:
    """Legacy shim — delegates to extract_with_deepseek()."""
    return extract_with_deepseek(text, api_key=api_key, model=model, workers=workers)
