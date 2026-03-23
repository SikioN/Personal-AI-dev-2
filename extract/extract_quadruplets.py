#!/usr/bin/env python3
"""
extract_quadruplets.py
======================
Извлекает квадруплеты (s, r, o, t) из документов (PDF, DOCX, PPTX),
верифицирует сущности и свойства через Wikidata API,
сохраняет результат в JSONL/JSON.

Зависимости:
    pip install ollama openai pymupdf python-docx python-pptx requests tqdm

Использование:
    # С Ollama (локально):
    python extract_quadruplets.py --input_dir new_docs --output out.jsonl

    # С OpenAI-совместимым API:
    python extract_quadruplets.py --input_dir new_docs --output out.jsonl \
        --api openai --api_key sk-... --model gpt-4o


Зависимости:
    pip install ollama openai pymupdf python-docx python-pptx requests tqdm

Использование:
    # С Ollama (локально):
    python extract_quadruplets.py --input_dir new_docs --output out.jsonl

    # С DeepSeek:
    python extract_quadruplets.py --input_dir new_docs --output out.jsonl \\
        --api deepseek --api_key sk-...

    # С OpenAI-совместимым API:
    python extract_quadruplets.py --input_dir new_docs --output out.jsonl \\
        --api openai --api_key sk-... --model gpt-4o
"""

import os
import re
import sys
import json
import time
import argparse
import logging
from pathlib import Path
from typing import Optional
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime, timezone

import requests
from tqdm import tqdm

logging.basicConfig(level=logging.INFO, format="%(levelname)s | %(message)s")
log = logging.getLogger(__name__)

# ─────────────────────────────────────────────
# СТАТИСТИКА ТОКЕНОВ
# ─────────────────────────────────────────────

_token_stats: dict[str, int] = {"prompt": 0, "completion": 0, "total": 0, "calls": 0}


def _add_tokens(prompt: int = 0, completion: int = 0) -> None:
    _token_stats["prompt"] += prompt
    _token_stats["completion"] += completion
    _token_stats["total"] += prompt + completion
    _token_stats["calls"] += 1


def _log_token_summary(prefix: str = "") -> None:
    log.info(
        f"{prefix}Tokens — total: {_token_stats['total']:,} "
        f"(prompt: {_token_stats['prompt']:,}, "
        f"completion: {_token_stats['completion']:,}), "
        f"LLM calls: {_token_stats['calls']}"
    )

# ─────────────────────────────────────────────
# 1. ЧТЕНИЕ ДОКУМЕНТОВ
# ─────────────────────────────────────────────

def read_pdf(path: str) -> str:
    import fitz  # PyMuPDF
    doc = fitz.open(path)
    return "\n".join(page.get_text() for page in doc)

def read_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

def read_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n".join(parts)

READERS = {
    ".pdf":  read_pdf,
    ".docx": read_docx,
    ".pptx": read_pptx,
    ".txt":  lambda p: Path(p).read_text(encoding="utf-8", errors="ignore"),
}

def read_document(path: str) -> Optional[str]:
    ext = Path(path).suffix.lower()
    reader = READERS.get(ext)
    if reader is None:
        log.warning(f"Unsupported format: {path}")
        return None
    try:
        text = reader(path)
        log.info(f"Read {path} → {len(text)} chars")
        return text
    except Exception as e:
        log.error(f"Cannot read {path}: {e}")
        return None

def collect_documents(input_dir: str) -> list[tuple[str, str]]:
    """Returns list of (filepath, text)."""
    docs = []
    for p in Path(input_dir).rglob("*"):
        if p.suffix.lower() in READERS and p.is_file():
            text = read_document(str(p))
            if text:
                docs.append((str(p), text))
    log.info(f"Collected {len(docs)} documents from '{input_dir}'")
    return docs

# ─────────────────────────────────────────────
# 2. КЭШ СВОЙСТВ WIKIDATA (P-идентификаторы)
# ─────────────────────────────────────────────

WIKIDATA_SPARQL = "https://query.wikidata.org/sparql"
WIKIDATA_ENTITY_SEARCH = "https://www.wikidata.org/w/api.php"

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "extract_data"

PROP_CACHE_FILE          = str(OUTPUT_DIR / "wikidata_properties_cache.json")

def fetch_all_wikidata_properties(use_cache: bool = True) -> dict[str, str]:
    """
    Возвращает словарь {label_lower: P-id} для популярных свойств Wikidata.
    При use_cache=True читает/пишет локальный файл кэша.
    """
    if use_cache and os.path.exists(PROP_CACHE_FILE):
        with open(PROP_CACHE_FILE, encoding="utf-8") as f:
            props = json.load(f)
        log.info(f"Loaded {len(props)} properties from cache '{PROP_CACHE_FILE}'")
        return props

    log.info("Fetching Wikidata properties via SPARQL (may take ~30s)…")
    query = """
    SELECT ?prop ?propLabel WHERE {
      ?prop wikibase:directClaim ?claim .
      SERVICE wikibase:label { bd:serviceParam wikibase:language "en". }
    }
    LIMIT 10000
    """
    headers = {"Accept": "application/sparql-results+json",
               "User-Agent": "QuadrupletExtractor/1.0"}
    resp = requests.get(WIKIDATA_SPARQL,
                        params={"query": query, "format": "json"},
                        headers=headers, timeout=120)
    resp.raise_for_status()
    data = resp.json()
    props: dict[str, str] = {}
    for row in data["results"]["bindings"]:
        pid = row["prop"]["value"].split("/")[-1]
        label = row["propLabel"]["value"].strip().lower()
        if pid.startswith("P") and label:
            props[label] = pid

    if use_cache:
        OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
        with open(PROP_CACHE_FILE, "w", encoding="utf-8") as f:
            json.dump(props, f, ensure_ascii=False, indent=2)
        log.info(f"Saved {len(props)} properties → '{PROP_CACHE_FILE}'")

    return props

def find_property(name: str, prop_map: dict[str, str]) -> Optional[str]:
    """Нечёткий поиск P-id по имени отношения."""
    key = name.strip().lower()
    if key in prop_map:
        return prop_map[key]
    # Частичное совпадение
    for label, pid in prop_map.items():
        if key in label or label in key:
            return pid
    return None


# ─────────────────────────────────────────────
# 2b. LLM-МАТЧИНГ СВОЙСТВ
# ─────────────────────────────────────────────

# Кэш: {relation_name_lower -> (wikidata_label, P-id)}
_prop_llm_cache: dict[str, tuple[str, str]] = {}
# Кэш результатов LLM-маппинга сохраняется на диск, чтобы не тратить токены повторно
PROP_LLM_CACHE_FILE      = str(OUTPUT_DIR / "wikidata_prop_llm_cache.json")

def _load_prop_llm_cache():
    global _prop_llm_cache
    if os.path.exists(PROP_LLM_CACHE_FILE):
        with open(PROP_LLM_CACHE_FILE, encoding="utf-8") as f:
            _prop_llm_cache = json.load(f)
        log.info(f"Loaded {len(_prop_llm_cache)} LLM property mappings from cache")

def _save_prop_llm_cache():
    with open(PROP_LLM_CACHE_FILE, "w", encoding="utf-8") as f:
        json.dump(_prop_llm_cache, f, ensure_ascii=False, indent=2)

def _get_top_candidates(query: str, prop_map: dict[str, str], top_n: int = 20) -> list[tuple[str, str]]:
    """
    Быстрый отбор топ-N кандидатов из prop_map по нескольким эвристикам:
    1. Слова запроса встречаются в лейбле (или наоборот)
    2. Сортировка по количеству совпадающих токенов
    """
    query_tokens = set(query.lower().split())
    scored: list[tuple[int, str, str]] = []

    for label, pid in prop_map.items():
        label_tokens = set(label.split())
        # Количество общих токенов
        common = len(query_tokens & label_tokens)
        if common > 0:
            scored.append((common, label, pid))

    # Сортируем по убыванию совпадений
    scored.sort(key=lambda x: -x[0])
    return [(label, pid) for _, label, pid in scored[:top_n]]


def find_property_with_llm(
    name: str,
    prop_map: dict[str, str],
    llm_caller,          # callable(prompt: str) -> str
    top_n: int = 20,
) -> tuple[Optional[str], Optional[str]]:
    """
    Трёхуровневый поиск свойства Wikidata:
      1. Точное совпадение (без LLM)
      2. Частичное совпадение по подстроке (без LLM)
      3. LLM выбирает лучший вариант из топ-N кандидатов

    Возвращает (matched_label, P-id) или (None, None).
    LLM вызывается только при необходимости, результат кэшируется на диск.
    """
    key = name.strip().lower()

    # ── Уровень 1: точное совпадение ──────────────────────────────────────
    if key in prop_map:
        log.debug(f"Property '{name}' → exact match → {prop_map[key]}")
        return key, prop_map[key]

    # ── Уровень 2: подстрочное совпадение ─────────────────────────────────
    for label, pid in prop_map.items():
        if key in label or label in key:
            log.debug(f"Property '{name}' → substring match '{label}' → {pid}")
            return label, pid

    # ── Уровень 3: LLM ────────────────────────────────────────────────────
    if key in _prop_llm_cache:
        cached = _prop_llm_cache[key]
        log.debug(f"Property '{name}' → LLM cache hit → {cached}")
        return cached[0], cached[1]

    candidates = _get_top_candidates(key, prop_map, top_n=top_n)

    if not candidates:
        log.warning(f"Property '{name}': no candidates found even for LLM")
        return None, None

    candidates_text = "\n".join(
        f"  - \"{label}\" ({pid})" for label, pid in candidates
    )

    prompt = (
        f"You are a Wikidata expert. A knowledge extraction system produced the relation label:\n"
        f"  \"{name}\"\n\n"
        f"Choose the SINGLE best matching Wikidata property from the list below.\n"
        f"Consider synonyms and paraphrases (e.g. 'received award' ≈ 'award received').\n\n"
        f"Candidates:\n{candidates_text}\n\n"
        f"Rules:\n"
        f"1. Reply ONLY with a JSON object, no extra text.\n"
        f"2. Format: {{\"label\": \"<wikidata label>\", \"id\": \"<P-id>\"}}\n"
        f"3. If none of the candidates is a good match, return: {{\"label\": null, \"id\": null}}\n"
    )

    try:
        raw_answer = llm_caller(prompt)
        raw_answer = re.sub(r"```(?:json)?", "", raw_answer).replace("```", "").strip()
        # Извлечь JSON из ответа
        match = re.search(r"\{[^{}]+\}", raw_answer, re.DOTALL)
        if not match:
            raise ValueError(f"No JSON object in LLM answer: {raw_answer[:200]}")
        parsed = json.loads(match.group(0))
        matched_label = parsed.get("label")
        matched_id    = parsed.get("id")

        if matched_label and matched_id:
            log.info(f"Property '{name}' → LLM matched '{matched_label}' ({matched_id})")
            _prop_llm_cache[key] = (matched_label, matched_id)
            _save_prop_llm_cache()
            return matched_label, matched_id
        else:
            log.warning(f"Property '{name}': LLM found no suitable match")
            _prop_llm_cache[key] = (None, None)
            _save_prop_llm_cache()
            return None, None

    except Exception as e:
        log.error(f"LLM property matching failed for '{name}': {e}")
        return None, None

# ─────────────────────────────────────────────
# 3. ВЕРИФИКАЦИЯ СУЩНОСТЕЙ В WIKIDATA
# ─────────────────────────────────────────────

_entity_cache: dict[str, Optional[str]] = {}

def search_entity_wikidata(name: str) -> Optional[str]:
    """Возвращает Q-идентификатор первого совпадения в Wikidata."""
    if name in _entity_cache:
        return _entity_cache[name]
    try:
        params = {
            "action": "wbsearchentities",
            "search": name,
            "language": "en",
            "limit": 1,
            "format": "json",
        }
        headers = {"User-Agent": "QuadrupletExtractor/1.0"}
        resp = requests.get(WIKIDATA_ENTITY_SEARCH, params=params,
                            headers=headers, timeout=10)
        resp.raise_for_status()
        results = resp.json().get("search", [])
        qid = results[0]["id"] if results else None
        _entity_cache[name] = qid
        time.sleep(0.1)  # вежливая задержка
        return qid
    except Exception as e:
        log.warning(f"Wikidata search failed for '{name}': {e}")
        _entity_cache[name] = None
        return None

# ─────────────────────────────────────────────
# 3b. РЕЕСТР КАСТОМНЫХ ID
# ─────────────────────────────────────────────
#
# Если сущность или свойство не найдены в Wikidata, назначаем им
# стабильный кастомный идентификатор:
#   CQ_000001, CQ_000002, … — для сущностей (аналог Q-id)
#   CP_000001, CP_000002, … — для свойств   (аналог P-id)
#
# Реестр сохраняется на диск → повторный анализ нового документа
# даст тем же именам те же ID.
#
# Структура файла custom_ids_registry.json:
# {
#   "entities": {"Apple Inc.": "CQ_000001", ...},
#   "properties": {"received award": "CP_000001", ...},
#   "next_entity_seq": 2,
#   "next_prop_seq": 2
# }

CUSTOM_ID_REGISTRY_FILE  = str(OUTPUT_DIR / "custom_ids_registry.json")
PROCESSED_MANIFEST_FILE  = str(OUTPUT_DIR / "processed_files.json")

_custom_registry: dict = {
    "entities":        {},   # name_lower → CQ_xxxxxx
    "properties":      {},   # name_lower → CP_xxxxxx
    "next_entity_seq": 1,
    "next_prop_seq":   1,
}


def _load_custom_registry() -> None:
    global _custom_registry
    if os.path.exists(CUSTOM_ID_REGISTRY_FILE):
        with open(CUSTOM_ID_REGISTRY_FILE, encoding="utf-8") as f:
            _custom_registry = json.load(f)
        log.info(
            f"Custom ID registry loaded: "
            f"{len(_custom_registry['entities'])} entities, "
            f"{len(_custom_registry['properties'])} properties"
        )


def _save_custom_registry() -> None:
    with open(CUSTOM_ID_REGISTRY_FILE, "w", encoding="utf-8") as f:
        json.dump(_custom_registry, f, ensure_ascii=False, indent=2)


# ─────────────────────────────────────────────
# МАНИФЕСТ ОБРАБОТАННЫХ ФАЙЛОВ
# ─────────────────────────────────────────────

import hashlib

_processed_manifest: dict[str, dict] = {}   # abs_path → {hash, processed_at, quads}


def _file_sha256(filepath: str) -> str:
    """SHA-256 первых 4 МБ файла — быстрая идентификация без полного чтения."""
    h = hashlib.sha256()
    with open(filepath, "rb") as f:
        h.update(f.read(4 * 1024 * 1024))
    return h.hexdigest()


def _load_processed_manifest() -> None:
    global _processed_manifest
    if os.path.exists(PROCESSED_MANIFEST_FILE):
        with open(PROCESSED_MANIFEST_FILE, encoding="utf-8") as f:
            _processed_manifest = json.load(f)
        log.info(f"Manifest loaded: {len(_processed_manifest)} already-processed file(s)")


def _save_processed_manifest() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    with open(PROCESSED_MANIFEST_FILE, "w", encoding="utf-8") as f:
        json.dump(_processed_manifest, f, ensure_ascii=False, indent=2)


def _is_processed(filepath: str) -> bool:
    """Возвращает True, если файл уже обработан (тот же хэш)."""
    key = str(Path(filepath).resolve())
    if key not in _processed_manifest:
        return False
    stored_hash = _processed_manifest[key].get("hash", "")
    return stored_hash == _file_sha256(filepath)


def _mark_processed(filepath: str, quads_count: int) -> None:
    key = str(Path(filepath).resolve())
    _processed_manifest[key] = {
        "hash": _file_sha256(filepath),
        "processed_at": datetime.now(timezone.utc).isoformat(timespec="seconds"),
        "quads": quads_count,
    }
    _save_processed_manifest()


def get_or_create_entity_id(name: str) -> str:
    """
    Возвращает существующий CQ-id для сущности или создаёт новый.
    Ключ — имя в нижнем регистре (регистронезависимо).
    """
    key = name.strip().lower()
    if key in _custom_registry["entities"]:
        return _custom_registry["entities"][key]

    seq = _custom_registry["next_entity_seq"]
    new_id = f"CQ_{seq:06d}"
    _custom_registry["entities"][key] = new_id
    _custom_registry["next_entity_seq"] = seq + 1
    _save_custom_registry()
    log.info(f"  [CustomID] New entity  '{name}' → {new_id}")
    return new_id


def get_or_create_property_id(name: str) -> str:
    """
    Возвращает существующий CP-id для свойства или создаёт новый.
    """
    key = name.strip().lower()
    if key in _custom_registry["properties"]:
        return _custom_registry["properties"][key]

    seq = _custom_registry["next_prop_seq"]
    new_id = f"CP_{seq:06d}"
    _custom_registry["properties"][key] = new_id
    _custom_registry["next_prop_seq"] = seq + 1
    _save_custom_registry()
    log.info(f"  [CustomID] New property '{name}' → {new_id}")
    return new_id


# ─────────────────────────────────────────────
# 4. LLM-БЭКЕНДЫ
# ─────────────────────────────────────────────

EXTRACTION_SYSTEM = """You are a knowledge graph construction expert.
Extract factual quadruplets (subject, relation, object, time) from the text.

Rules:
1. Subject and Object must be named entities (persons, organizations, places, works, concepts).
2. Relation must be a Wikidata property label (e.g. "spouse", "award received", "position held", "educated at", "employer", "member of", "country of citizenship", "has part" etc.).
3. Time: use years only. If exact year unknown, omit. Format: {"start": "YYYY", "end": "YYYY"}.
   If a single year, set start == end. If open-ended (still active), use current year as end.
4. Output ONLY a valid JSON object with a single key "quadruplets" containing the array. Do not output anything else.

Output format (strict JSON object):
{
  "quadruplets": [
    {
      "s": "Subject Name",
      "r": "relation label",
      "o": "Object Name",
      "t_start": "YYYY",
      "t_end": "YYYY"
    }
  ]
}

If no quadruplets found, return: {"quadruplets": []}
"""

def chunk_text(text: str, max_chars: int = 6000) -> list[str]:
    """Split text into overlapping chunks."""
    chunks = []
    step = max_chars - 500
    for i in range(0, len(text), step):
        chunks.append(text[i:i + max_chars])
    return chunks

def extract_with_ollama(text: str, model: str = "qwen3:8b") -> list[dict]:
    try:
        import ollama
    except ImportError:
        log.error("Install ollama: pip install ollama")
        return []

    client = ollama.Client(timeout=3000)
    chunks = chunk_text(text)
    all_quads = []
    bar = tqdm(chunks, desc="  chunks", unit="chunk", leave=False)
    for chunk in bar:
        try:
            response = client.chat(
                model=model,
                messages=[
                    {"role": "system", "content": EXTRACTION_SYSTEM},
                    {"role": "user", "content": f"Extract quadruplets from this text:\n\n{chunk}"},
                ],
            )
            raw = response.message.content
            # Ollama возвращает статистику токенов
            prompt_tok = getattr(response, "prompt_eval_count", 0) or 0
            compl_tok  = getattr(response, "eval_count", 0) or 0
            _add_tokens(prompt_tok, compl_tok)
            bar.set_postfix(tokens=_token_stats["total"])
            quads = parse_llm_json(raw)
            all_quads.extend(quads)
        except Exception as e:
            log.error(f"Ollama error: {e}")
    return all_quads

def _process_chunk_openai(chunk: str, client, model: str) -> tuple[list[dict], int, int]:
    """Обрабатывает один чанк через OpenAI-совместимый API. Возвращает (quads, prompt_tok, compl_tok)."""
    use_json_mode = any(x in model for x in ("gpt-4", "deepseek"))
    resp = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": EXTRACTION_SYSTEM},
            {"role": "user", "content": f"Extract quadruplets from this text:\n\n{chunk}"},
        ],
        temperature=0,
        response_format={"type": "json_object"} if use_json_mode else None,
    )
    raw = resp.choices[0].message.content
    usage = resp.usage
    prompt_tok = usage.prompt_tokens if usage else 0
    compl_tok  = usage.completion_tokens if usage else 0
    return parse_llm_json(raw), prompt_tok, compl_tok


def extract_with_openai(text: str, api_key: str, model: str = "gpt-4o",
                        base_url: Optional[str] = None,
                        workers: int = 4) -> tuple[list[dict], int]:
    try:
        from openai import OpenAI
    except ImportError:
        log.error("Install openai: pip install openai")
        return [], 0

    client = OpenAI(api_key=api_key, base_url=base_url)
    chunks = chunk_text(text)
    all_quads: list[dict] = []

    failed_chunks = 0
    with ThreadPoolExecutor(max_workers=workers) as executor:
        futures = {executor.submit(_process_chunk_openai, chunk, client, model): i
                   for i, chunk in enumerate(chunks)}
        bar = tqdm(as_completed(futures), total=len(chunks),
                   desc="  chunks", unit="chunk", leave=False)
        for future in bar:
            try:
                quads, prompt_tok, compl_tok = future.result()
                all_quads.extend(quads)
                _add_tokens(prompt_tok, compl_tok)
                bar.set_postfix(tokens=f"{_token_stats['total']:,}")
            except Exception as e:
                log.error(f"OpenAI chunk error: {e}")
                failed_chunks += 1

    if failed_chunks:
        log.warning(f"extract_with_openai: {failed_chunks}/{len(chunks)} chunks failed")
    return all_quads, failed_chunks

def parse_llm_json(raw: str) -> list[dict]:
    """Extract JSON array from LLM response (handles markdown fences)."""
    raw = raw.strip()
    # Remove markdown fences
    raw = re.sub(r"```(?:json)?", "", raw).replace("```", "").strip()
    # Find JSON array
    match = re.search(r"\[.*\]", raw, re.DOTALL)
    if not match:
        # Maybe it's wrapped in {"quadruplets": [...]}
        match = re.search(r'\{.*"(?:quadruplets|data|results)":\s*(\[.*\])', raw, re.DOTALL)
        if match:
            raw = match.group(1)
        else:
            log.warning(f"No JSON array found in LLM response: {raw[:200]}")
            return []
    else:
        raw = match.group(0)
    try:
        return json.loads(raw)
    except json.JSONDecodeError as e:
        log.warning(f"JSON parse error: {e}\nRaw: {raw[:300]}")
        return []

# ─────────────────────────────────────────────
# 4b. LLM-CALLER — единый вызов для любого бэкенда
# ─────────────────────────────────────────────

DEEPSEEK_BASE_URL = "https://api.deepseek.com"
DEEPSEEK_DEFAULT_MODEL = "deepseek-chat"
DEEPSEEK_API_KEY = os.environ.get("DEEPSEEK_API_KEY", "")


def make_llm_caller(api: str, model: str, api_key: str = "",
                    base_url: Optional[str] = None):
    """
    Возвращает callable(prompt: str) -> str для лёгких одиночных запросов
    (в частности, для LLM-матчинга свойств).
    """
    if api == "ollama":
        try:
            import ollama as _ollama
        except ImportError:
            log.error("Install ollama: pip install ollama")
            return lambda p: "{}"

        def call_ollama(prompt: str) -> str:
            resp = _ollama.chat(
                model=model,
                messages=[{"role": "user", "content": prompt}],
            )
            prompt_tok = getattr(resp, "prompt_eval_count", 0) or 0
            compl_tok  = getattr(resp, "eval_count", 0) or 0
            _add_tokens(prompt_tok, compl_tok)
            return resp["message"]["content"]

        return call_ollama

    else:  # openai / deepseek / compatible
        try:
            from openai import OpenAI as _OAI
        except ImportError:
            log.error("Install openai: pip install openai")
            return lambda p: "{}"

        client = _OAI(api_key=api_key or "none", base_url=base_url)

        def call_openai(prompt: str) -> str:
            resp = client.chat.completions.create(
                model=model,
                messages=[{"role": "user", "content": prompt}],
                temperature=0,
            )
            if resp.usage:
                _add_tokens(resp.usage.prompt_tokens, resp.usage.completion_tokens)
            return resp.choices[0].message.content

        return call_openai


# ─────────────────────────────────────────────
# 5. ПОСТРОЕНИЕ ФИНАЛЬНОГО КВАДРУПЛЕТА
# ─────────────────────────────────────────────

def build_quadruplet(raw: dict, prop_map: dict[str, str],
                     llm_caller=None) -> Optional[dict]:
    """
    Принимает сырой dict от LLM, верифицирует через Wikidata,
    возвращает квадруплет в целевом формате или None при ошибке.
    """
    s_name = str(raw.get("s", "")).strip()
    r_name = str(raw.get("r", "")).strip()
    o_name = str(raw.get("o", "")).strip()
    t_start = str(raw.get("t_start", "")).strip()
    t_end   = str(raw.get("t_end", t_start)).strip()

    if not s_name or not r_name or not o_name:
        return None

    # Верификация сущностей: Wikidata → иначе кастомный CQ-id
    s_id = search_entity_wikidata(s_name) or get_or_create_entity_id(s_name)
    o_id = search_entity_wikidata(o_name) or get_or_create_entity_id(o_name)

    # Верификация свойства: точный/частичный → LLM → иначе кастомный CP-id
    if llm_caller is not None:
        matched_label, r_id = find_property_with_llm(r_name, prop_map, llm_caller)
        display_r_name = matched_label if matched_label else r_name
    else:
        r_id = find_property(r_name, prop_map)
        display_r_name = r_name

    if not r_id:
        r_id = get_or_create_property_id(display_r_name)

    # Время
    t_name = f"{t_start} - {t_end}" if t_start and t_end else "unknown"
    t_block = {
        "name": t_name,
        "type": "time",
        "prop": {"start": t_start, "end": t_end},
    } if t_start else {"name": "unknown", "type": "time", "prop": {}}

    quad = {
        "s": {"name": s_name, "id": s_id},
        "r": {"name": display_r_name, "id": r_id, "type": "simple"},
        "o": {"name": o_name, "id": o_id},
        "t": t_block,
        # Служебное поле для отладки: что вернула LLM до нормализации
        "_r_raw": r_name if display_r_name != r_name else None,
    }
    # Убираем None-поля отладки
    if quad["_r_raw"] is None:
        del quad["_r_raw"]
    return quad

# ─────────────────────────────────────────────
# 6. ДЕДУПЛИКАЦИЯ
# ─────────────────────────────────────────────

def deduplicate(quads: list[dict]) -> list[dict]:
    seen = set()
    result = []
    for q in quads:
        key = (q["s"]["name"], q["r"]["name"], q["o"]["name"],
               q["t"]["prop"].get("start"), q["t"]["prop"].get("end"))
        if key not in seen:
            seen.add(key)
            result.append(q)
    return result

# ─────────────────────────────────────────────
# 7. ГЛАВНАЯ ФУНКЦИЯ
# ─────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Extract knowledge graph quadruplets from documents")
    parser.add_argument("--input_dir", default="new_docs", help="Folder with documents")
    parser.add_argument("--output", default=str(OUTPUT_DIR / "quadruplets.jsonl"),
                        help="Output JSONL file")
    parser.add_argument("--api", choices=["ollama", "openai", "deepseek", "compatible"],
                        default="ollama", help="LLM backend")
    parser.add_argument("--model", default="", help="Model name (default depends on --api)")
    parser.add_argument("--api_key", default="", help="API key (for openai/deepseek/compatible)")
    parser.add_argument("--base_url", default=None, help="Base URL for compatible API")
    parser.add_argument("--workers", type=int, default=4,
                        help="Parallel chunk workers for OpenAI/DeepSeek (default: 4)")
    parser.add_argument("--no_prop_cache", action="store_true",
                        help="Force re-fetch Wikidata properties")
    parser.add_argument("--reprocess", action="store_true",
                        help="Ignore processed-files manifest and reprocess everything")
    parser.add_argument("--mark-done", action="store_true",
                        help="Mark all files in input_dir as already processed (no LLM calls)")
    args = parser.parse_args()

    # ── Режим пометки без обработки ──────────────────────────────
    if args.mark_done:
        OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
        _load_processed_manifest()
        files = list(collect_documents(args.input_dir))
        if not files:
            log.warning("No documents found in input_dir.")
            sys.exit(0)
        bar = tqdm(files, desc="Marking", unit="file")
        for filepath, _ in bar:
            bar.set_description(f"Marking: {Path(filepath).name[:50]}")
            _mark_processed(filepath, quads_count=-1)   # -1 = кол-во неизвестно
            log.info(f"  Marked: {filepath}")
        log.info(f"Done. {len(files)} file(s) added to manifest → {PROCESSED_MANIFEST_FILE}")
        sys.exit(0)

    # Настройки по умолчанию для DeepSeek
    if args.api == "deepseek":
        if not args.api_key:
            args.api_key = DEEPSEEK_API_KEY
        if not args.base_url:
            args.base_url = DEEPSEEK_BASE_URL
        if not args.model:
            args.model = DEEPSEEK_DEFAULT_MODEL
        log.info(f"Using DeepSeek backend: {args.base_url}, model={args.model}")
    elif not args.model:
        args.model = "qwen3:8b" if args.api == "ollama" else "gpt-4o"

    # Создаём выходную папку
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    # Шаг 1: загрузка свойств Wikidata
    prop_map = fetch_all_wikidata_properties(use_cache=not args.no_prop_cache)

    # Шаг 1b: загрузка кэша LLM-маппинга свойств
    _load_prop_llm_cache()

    # Шаг 1c: загрузка реестра кастомных ID
    _load_custom_registry()

    # Шаг 1e: загрузка манифеста обработанных файлов
    _load_processed_manifest()

    # Шаг 1d: создаём llm_caller для матчинга свойств
    llm_caller = make_llm_caller(
        api=args.api,
        model=args.model,
        api_key=args.api_key,
        base_url=args.base_url,
    )

    # Шаг 2: чтение документов
    if not os.path.exists(args.input_dir):
        log.error(f"Input directory '{args.input_dir}' not found.")
        sys.exit(1)
    documents = collect_documents(args.input_dir)

    if not documents:
        log.warning("No documents found!")
        sys.exit(0)

    # Фильтруем уже обработанные файлы
    new_documents = documents if args.reprocess else \
                    [(fp, txt) for fp, txt in documents if not _is_processed(fp)]
    skipped = len(documents) - len(new_documents)
    if skipped:
        log.info(f"Skipping {skipped} already-processed file(s). "
                 f"Use --reprocess to force reprocessing.")
    if not new_documents:
        log.info("All files already processed. Nothing to do.")
        sys.exit(0)

    # Шаг 3: извлечение квадруплетов
    all_quadruplets: list[dict] = []
    start_total = time.monotonic()

    doc_bar = tqdm(new_documents, desc="Documents", unit="doc")
    for doc_idx, (filepath, text) in enumerate(doc_bar):
        doc_start = time.monotonic()
        doc_bar.set_description(f"Doc {doc_idx+1}/{len(new_documents)}: {Path(filepath).name[:40]}")
        log.info(f"\n{'='*60}\nProcessing: {filepath}  ({len(text):,} chars)\n{'='*60}")

        if args.api == "ollama":
            raw_quads = extract_with_ollama(text, model=args.model)
        elif args.api in ("openai", "deepseek", "compatible"):
            raw_quads, _chunk_errors = extract_with_openai(
                text,
                api_key=args.api_key,
                model=args.model,
                base_url=args.base_url,
                workers=args.workers,
            )
        else:
            raw_quads = []

        doc_elapsed = time.monotonic() - doc_start
        log.info(f"LLM extracted {len(raw_quads)} raw quadruplets in {doc_elapsed:.1f}s")
        _log_token_summary(prefix="  ")

        # Прогноз оставшегося времени по документам
        elapsed_total = time.monotonic() - start_total
        docs_done = doc_idx + 1
        docs_left = len(new_documents) - docs_done
        if docs_done > 0 and docs_left > 0:
            eta_sec = (elapsed_total / docs_done) * docs_left
            m, s = divmod(int(eta_sec), 60)
            h, m = divmod(m, 60)
            log.info(f"  ETA for remaining {docs_left} doc(s): {h:02d}:{m:02d}:{s:02d}")

        # Верификация и построение финальных квадруплетов
        doc_quads: list[dict] = []
        for rq in raw_quads:
            quad = build_quadruplet(rq, prop_map, llm_caller=llm_caller)
            if quad:
                doc_quads.append(quad)
                print(json.dumps(quad, ensure_ascii=False))
        all_quadruplets.extend(doc_quads)

        # Отмечаем файл как обработанный
        _mark_processed(filepath, len(doc_quads))

    # Шаг 4: дедупликация
    all_quadruplets = deduplicate(all_quadruplets)
    log.info(f"\nTotal unique quadruplets: {len(all_quadruplets)}")

    # Итоговая статистика токенов
    _log_token_summary(prefix="FINAL ")
    total_sec = time.monotonic() - start_total
    h, rem = divmod(int(total_sec), 3600)
    m, s = divmod(rem, 60)
    log.info(f"Total time: {h:02d}:{m:02d}:{s:02d}")

    # Шаг 5: сохранение
    out_path = Path(args.output)

    # JSONL
    with open(out_path, "w", encoding="utf-8") as f:
        for q in all_quadruplets:
            f.write(json.dumps(q, ensure_ascii=False) + "\n")
    log.info(f"Saved JSONL → {out_path}")

    # JSON (полный массив)
    json_path = out_path.with_suffix(".json")
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(all_quadruplets, f, ensure_ascii=False, indent=2)
    log.info(f"Saved JSON  → {json_path}")

    print(f"\n✅ Done: {len(all_quadruplets)} quadruplets saved to {out_path} and {json_path}")

if __name__ == "__main__":
    main()
    